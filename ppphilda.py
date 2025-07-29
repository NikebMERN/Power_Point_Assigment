from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Slide content definitions
slides_data = [
    {
        "title": "Diagnostic Tools of ICT Assigment",
        "content": "Ensuring Optimal Performance and Troubleshooting\nNikodimos Elias\n12 of May 2024",
        "img": "slide1_bg.jpg"
    },
    {
        "title": "Introduction to ICT Diagnostic Tools",
        "content": (
            "What are ICT Diagnostic Tools?\n"
            "• Tools, techniques, and software used to identify, analyze, and resolve problems within ICT systems.\n\n"
            "Importance of ICT Diagnostic Tools:\n"
            "• Minimize downtime and disruptions\n"
            "• Improve system performance and reliability\n"
            "• Enhance security by identifying vulnerabilities\n"
            "• Facilitate efficient troubleshooting\n"
            "• Enable proactive maintenance"
        ),
        "img": "slide2_bg.jpg"
    },
    {
        "title": "Hardware Diagnostic Tools",
        "content": (
            "Internal Hardware Tools:\n"
            "• POST (Power-On Self-Test)\n"
            "• Multimeters\n"
            "• Logic Probes\n\n"
            "External Hardware Tools:\n"
            "• Cable testers\n"
            "• Hard drive testers\n"
            "• RAM testers"
        ),
        "img": "slide3_bg.jpg"
    },
    {
        "title": "Software Diagnostic Tools",
        "content": (
            "System Monitoring Tools:\n"
            "• OS Utilities (Task Manager, Resource Monitor)\n"
            "• Third-party system monitoring tools\n\n"
            "Network Analysis Tools:\n"
            "• ping\n"
            "• traceroute/tracert\n"
            "• Wireshark\n"
            "• nslookup/dig"
        ),
        "img": "slide4_bg.jpg"
    },
    {
        "title": "Examples and Applications of ICT Diagnostic Tools",
        "content": (
            "Network Troubleshooting:\n"
            "• ping, traceroute, Wireshark\n\n"
            "System Performance Optimization:\n"
            "• Monitor CPU, memory, disk I/O\n\n"
            "Security Diagnostics:\n"
            "• Vulnerability scanners, IDS/IPS\n\n"
            "Hardware Diagnostics:\n"
            "• Memtest86, SMART"
        ),
        "img": "slide5_bg.jpg"
    },
    {
        "title": "Emerging Trends and Conclusion",
        "content": (
            "Emerging Trends:\n"
            "• AI and ML for predictive diagnostics\n"
            "• Cloud-based remote diagnostics\n"
            "• IoT monitoring and analysis\n\n"
            "Conclusion:\n"
            "• Diagnostic tools ensure proactive problem-solving\n"
            "• Minimize downtime, enhance performance\n"
            "• Stay current with emerging trends"
        ),
        "img": "slide6_bg.jpg"
    },
]

# Function to add a new slide
def add_slide(prs, title, content, bg_img):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add background image
    slide.shapes.add_picture(bg_img, 0, 0, width=prs.slide_width, height=prs.slide_height)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(12), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Add body content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    for line in content.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)

# Add all slides
for slide in slides_data:
    add_slide(prs, slide["title"], slide["content"], slide["img"])

# Save the presentation
prs.save("Diagnostic_Tools_of_ICT_Presentation.pptx")
